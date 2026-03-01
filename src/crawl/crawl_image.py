import base64
import hashlib
import re
import subprocess
import threading
import time
import zipfile
from io import BytesIO
from pathlib import Path
from typing import Dict, Optional

try:
    from core.config import Settings
except ImportError:
    from src.core.config import Settings


def sanitize_filename(name: str) -> str:
    return re.sub(r'[\\/*?:"<>|]', "_", name or "")


class VisionAnalyzer:
    """
    Runtime router:
    - GPU + local model load success -> Qwen local
    - otherwise -> Groq Vision
    """

    def __init__(self):
        self.mode = (Settings.LLM_MODE or "auto").lower()
        self.local_model_id = Settings.LOCAL_VLM_ID
        self.groq_model = Settings.GROQ_VISION_MODEL
        self._local_ready = False
        self._groq_ready = False
        self._init_attempted = False

        self._model = None
        self._processor = None
        self._torch = None
        self._process_vision_info = None
        self._groq_client = None
        self._lock = threading.Lock()

    def _init_local(self) -> bool:
        if self._local_ready:
            return True
        try:
            import torch
            from PIL import Image
            from qwen_vl_utils import process_vision_info
            from transformers import AutoProcessor, Qwen2VLForConditionalGeneration
        except Exception:
            return False

        if not torch.cuda.is_available():
            return False

        try:
            model = Qwen2VLForConditionalGeneration.from_pretrained(
                self.local_model_id,
                torch_dtype=torch.bfloat16,
                device_map="auto",
            )
            processor = AutoProcessor.from_pretrained(
                self.local_model_id,
                min_pixels=256 * 28 * 28,
                max_pixels=1024 * 28 * 28,
            )
            self._model = model
            self._processor = processor
            self._torch = torch
            self._process_vision_info = process_vision_info
            self._local_ready = True
            return True
        except Exception:
            return False

    def _init_groq(self) -> bool:
        if self._groq_ready:
            return True
        if not Settings.GROQ_API_KEY:
            return False
        try:
            from groq import Groq
        except Exception:
            return False
        self._groq_client = Groq(api_key=Settings.GROQ_API_KEY)
        self._groq_ready = True
        return True

    def _init_once(self) -> None:
        if self._init_attempted:
            return
        self._init_attempted = True

        if self.mode == "qwen_local":
            self._local_ready = self._init_local()
            if not self._local_ready:
                self._groq_ready = self._init_groq()
            return

        if self.mode == "groq":
            self._groq_ready = self._init_groq()
            return

        # auto
        self._local_ready = self._init_local()
        if not self._local_ready:
            self._groq_ready = self._init_groq()

    @staticmethod
    def _bytes_to_data_url(image_bytes: BytesIO) -> str:
        raw = image_bytes.getvalue()
        b64 = base64.b64encode(raw).decode("utf-8")
        return f"data:image/jpeg;base64,{b64}"

    def _analyze_local(self, image_bytes: BytesIO, prompt: str) -> str:
        if not self._local_ready:
            return ""
        try:
            from PIL import Image

            image_bytes.seek(0)
            img = Image.open(image_bytes).convert("RGB")
            if img.width < 80 or img.height < 80:
                return ""

            messages = [
                {
                    "role": "user",
                    "content": [
                        {"type": "image", "image": img},
                        {"type": "text", "text": prompt},
                    ],
                }
            ]

            text = self._processor.apply_chat_template(
                messages, tokenize=False, add_generation_prompt=True
            )
            image_inputs, video_inputs = self._process_vision_info(messages)
            inputs = self._processor(
                text=[text],
                images=image_inputs,
                videos=video_inputs,
                padding=True,
                return_tensors="pt",
            ).to(self._model.device)

            with self._lock:
                generated = self._model.generate(**inputs, max_new_tokens=768)

            trimmed = [
                out_ids[len(in_ids) :]
                for in_ids, out_ids in zip(inputs.input_ids, generated)
            ]
            out = self._processor.batch_decode(trimmed, skip_special_tokens=True)[0]

            del inputs, generated, trimmed
            if self._torch and self._torch.cuda.is_available():
                self._torch.cuda.empty_cache()
            return out or ""
        except Exception:
            return ""

    def _analyze_groq(self, image_src: str, prompt: str) -> str:
        if not self._groq_ready:
            return ""
        if image_src.startswith("data:image") and len(image_src) < 1000:
            return ""
        try:
            completion = self._groq_client.chat.completions.create(
                model=self.groq_model,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": prompt},
                            {"type": "image_url", "image_url": {"url": image_src}},
                        ],
                    }
                ],
                temperature=0.0,
                max_tokens=1200,
            )
            return completion.choices[0].message.content or ""
        except Exception:
            return ""

    def analyze_image(self, image_src: str, prompt: str) -> str:
        self._init_once()
        if self._groq_ready:
            return self._analyze_groq(image_src, prompt)
        return ""

    def analyze_bytes(self, image_bytes: BytesIO, prompt: str) -> str:
        self._init_once()
        if self._local_ready:
            out = self._analyze_local(image_bytes, prompt)
            if out:
                return out
        if self._groq_ready:
            data_url = self._bytes_to_data_url(image_bytes)
            return self._analyze_groq(data_url, prompt)
        return ""


VISION = VisionAnalyzer()


def analyze_image(image_src: str) -> str:
    prompt = (
        "Extract all visible text and key structured details from this image. "
        "Return concise markdown only."
    )
    return VISION.analyze_image(image_src=image_src, prompt=prompt)


def analyze_image_from_memory(image_bytes: BytesIO, alt_text: str = "") -> str:
    prompt = (
        "Extract all visible text and key structured details from this image. "
        "Return concise markdown only."
    )
    if alt_text:
        prompt += f"\nAdditional hint: {alt_text[:120]}"
    return VISION.analyze_bytes(image_bytes=image_bytes, prompt=prompt)


def _download_file(session, url, save_path: Path, referer=None):
    save_path.parent.mkdir(parents=True, exist_ok=True)
    headers = {"Referer": referer} if referer else {}

    try:
        resp = session.get(url, headers=headers, verify=False, timeout=30, stream=True)
        resp.raise_for_status()
        h = hashlib.sha256()
        size = 0
        with open(save_path, "wb") as f:
            for chunk in resp.iter_content(chunk_size=8192):
                if not chunk:
                    continue
                f.write(chunk)
                h.update(chunk)
                size += len(chunk)
        return {"size": size, "sha256": h.hexdigest(), "status": "success"}
    except Exception as e:
        return {"status": "error", "error": str(e)}


def _download_image_to_memory(session, url, referer=None):
    headers = {"Referer": referer} if referer else {}
    try:
        resp = session.get(url, headers=headers, verify=False, timeout=20, stream=True)
        resp.raise_for_status()
        image_bytes = BytesIO()
        for chunk in resp.iter_content(chunk_size=8192):
            if chunk:
                image_bytes.write(chunk)
        image_bytes.seek(0)
        return image_bytes
    except Exception:
        return None


def _get_module_version(module_name: str) -> str:
    try:
        import importlib.metadata as importlib_metadata

        return importlib_metadata.version(module_name)
    except Exception:
        return "unknown"


def _score_confidence(text: str, parse_error: str = "") -> float:
    if parse_error and not text:
        return 0.0
    text_len = len((text or "").strip())
    if text_len >= 400:
        return 0.95
    if text_len >= 150:
        return 0.85
    if text_len >= 40:
        return 0.7
    if text_len > 0:
        return 0.5
    return 0.0


def extract_text_with_meta(file_path: Path, ext: str) -> Dict[str, object]:
    text = ""
    max_chars = 20000
    parser_name = "none"
    parser_version = "unknown"
    parse_error = ""
    extraction_method = "text_parser"

    def _clean_text(raw: str) -> str:
        if not raw:
            return ""
        cleaned = re.sub(r"\r\n?", "\n", raw)
        cleaned = re.sub(r"[ \t]+", " ", cleaned)
        cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
        return cleaned.strip()[:max_chars]

    ext = ext.lower()
    try:
        if ext == ".pdf":
            import fitz

            fitz.TOOLS.mupdf_display_errors(False)
            with fitz.open(str(file_path)) as doc:
                texts = [page.get_text().strip() for page in doc if page.get_text().strip()]
                text = _clean_text("\n\n".join(texts))
            parser_name = "pymupdf"
            parser_version = _get_module_version("PyMuPDF")

        elif ext == ".docx":
            import docx
            from docx.oxml.table import CT_Tbl
            from docx.oxml.text.paragraph import CT_P
            from docx.table import Table
            from docx.text.paragraph import Paragraph

            doc = docx.Document(str(file_path))
            texts = []
            for element in doc.element.body:
                if isinstance(element, CT_P):
                    p = Paragraph(element, doc)
                    if p.text.strip():
                        texts.append(p.text.strip())
                elif isinstance(element, CT_Tbl):
                    table = Table(element, doc)
                    rows = []
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                        if row_text:
                            rows.append(row_text)
                    if rows:
                        texts.append("\n[표]\n" + "\n".join(rows))
            text = _clean_text("\n\n".join(texts))
            parser_name = "python-docx"
            parser_version = _get_module_version("python-docx")

        elif ext == ".hwpx":
            try:
                import hwpx  # type: ignore

                parser_name = "python-hwpx"
                parser_version = _get_module_version("python-hwpx")
                doc_obj = None
                if hasattr(hwpx, "HWPX"):
                    doc_obj = hwpx.HWPX(str(file_path))
                elif hasattr(hwpx, "open"):
                    doc_obj = hwpx.open(str(file_path))

                if doc_obj is not None:
                    for attr in ("text", "content"):
                        if hasattr(doc_obj, attr):
                            value = getattr(doc_obj, attr)
                            if isinstance(value, str) and value.strip():
                                text = _clean_text(value)
                                break
                    if not text:
                        for method_name in ("get_text", "to_text", "extract_text"):
                            if hasattr(doc_obj, method_name):
                                value = getattr(doc_obj, method_name)()
                                if isinstance(value, str) and value.strip():
                                    text = _clean_text(value)
                                    break
            except Exception:
                pass

            if not text:
                from lxml import etree

                parser_name = "hwpx-xml-fallback"
                parser_version = _get_module_version("lxml")
                with zipfile.ZipFile(file_path, "r") as zf:
                    xml_candidates = [
                        name
                        for name in zf.namelist()
                        if name.endswith(".xml") and ("Contents" in name or "section" in name.lower())
                    ]
                    chunks = []
                    for xml_name in xml_candidates[:30]:
                        with zf.open(xml_name) as fp:
                            tree = etree.parse(fp)
                            nodes = tree.xpath("//*[local-name()='t']/text()")
                            if nodes:
                                chunks.append("\n".join([n.strip() for n in nodes if n and n.strip()]))
                    text = _clean_text("\n\n".join(chunks))

        elif ext == ".hwp":
            try:
                result = subprocess.run(
                    ["hwp5txt", str(file_path)],
                    capture_output=True,
                    text=True,
                    encoding="utf-8",
                    timeout=12,
                )
                if result.returncode == 0 and result.stdout.strip():
                    text = _clean_text(result.stdout.strip())
                    parser_name = "pyhwp-hwp5txt"
                    parser_version = _get_module_version("pyhwp")
                else:
                    raise RuntimeError("hwp5txt empty")
            except Exception:
                try:
                    result = subprocess.run(
                        ["hwp-extract", str(file_path)],
                        capture_output=True,
                        text=True,
                        encoding="utf-8",
                        timeout=15,
                    )
                    if result.returncode == 0 and result.stdout.strip():
                        text = _clean_text(result.stdout.strip())
                        parser_name = "hwp-extract"
                        parser_version = _get_module_version("hwp-extract")
                except Exception:
                    pass

            if not text:
                import olefile

                if olefile.isOleFile(str(file_path)):
                    parser_name = "olefile-fallback"
                    parser_version = _get_module_version("olefile")
                    ole = olefile.OleFileIO(str(file_path))
                    texts = []
                    for entry in ole.listdir():
                        if entry and entry[0] == "BodyText":
                            stream = ole.openstream(entry)
                            data = stream.read()
                            decoded = data.decode("utf-16le", errors="ignore")
                            cleaned = "".join([c for c in decoded if c.isprintable() or c in "\n\t "]).strip()
                            if cleaned:
                                texts.append(cleaned)
                    ole.close()
                    text = _clean_text("\n\n".join(texts))

        elif ext in [".xlsx", ".xls"]:
            import pandas as pd

            all_sheets = pd.read_excel(str(file_path), sheet_name=None, dtype=str)
            chunks = []
            for sheet_name, df in all_sheets.items():
                if df is None or df.empty:
                    continue
                df = df.fillna("")
                head = df.head(200)
                rows = []
                for _, row in head.iterrows():
                    vals = [str(v).strip() for v in row.tolist() if str(v).strip()]
                    if vals:
                        rows.append(" | ".join(vals))
                if rows:
                    chunks.append(f"[시트: {sheet_name}]\n" + "\n".join(rows))
            text = _clean_text("\n\n".join(chunks))
            parser_name = "pandas-excel"
            parser_version = _get_module_version("pandas")

        elif ext == ".pptx":
            from pptx import Presentation

            prs = Presentation(str(file_path))
            slides_text = []
            for i, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        t = str(shape.text).strip()
                        if t:
                            parts.append(t)
                if parts:
                    slides_text.append(f"[슬라이드 {i}]\n" + "\n".join(parts))
            text = _clean_text("\n\n".join(slides_text))
            parser_name = "python-pptx"
            parser_version = _get_module_version("python-pptx")

        elif ext in [".txt", ".csv"]:
            for enc in ["utf-8", "cp949", "euc-kr", "latin-1"]:
                try:
                    with open(file_path, "r", encoding=enc, errors="ignore") as f:
                        text = _clean_text(f.read())
                    if text:
                        parser_name = f"text-encoding:{enc}"
                        parser_version = "builtin"
                        break
                except Exception:
                    continue
    except Exception as exc:
        parse_error = str(exc)

    text = _clean_text(text)
    return {
        "text": text,
        "parser_name": parser_name,
        "parser_version": parser_version,
        "parse_confidence": _score_confidence(text, parse_error),
        "parse_error": parse_error,
        "extraction_method": extraction_method,
    }


def _extract_text_from_file(file_path: Path, ext: str) -> str:
    return str(extract_text_with_meta(file_path, ext).get("text", ""))
